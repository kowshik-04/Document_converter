from PIL import Image
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
from fpdf import FPDF

# Convert Image to PDF
def img_to_pdf(input_path, output_path):
    """Convert image to PDF."""
    img = Image.open(input_path)
    img.save(output_path, "PDF", resolution=100.0)

# Convert Image to Word
def img_to_word(input_path, output_path):
    """Convert image to Word document using OCR."""
    img = Image.open(input_path)
    text = pytesseract.image_to_string(img)
    doc = Document()
    doc.add_paragraph(text)
    doc.save(output_path)

# Convert PDF to Image
def pdf_to_image(input_path, output_path):
    """Convert PDF to images."""
    images = convert_from_path(input_path)
    for i, image in enumerate(images):
        image.save(f"{output_path}_page_{i+1}.jpg", 'JPEG')

# OCR for PDF and Image (convert to text)
def ocr_for_pdf(input_path):
    """Perform OCR on PDF to extract text."""
    text = ""
    images = convert_from_path(input_path)
    for image in images:
        text += pytesseract.image_to_string(image)
    return text

def ocr_for_image(input_path):
    """Perform OCR on image to extract text."""
    img = Image.open(input_path)
    return pytesseract.image_to_string(img)

def pdf_to_word(input_path, output_path):
    """Convert PDF to Word document."""
    reader = PdfReader(input_path)
    doc = Document()

    for page in reader.pages:
        text = page.extract_text()
        doc.add_paragraph(text)

    doc.save(output_path)

def word_to_pdf(input_path, output_path):
    """Convert Word to PDF."""
    doc = Document(input_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for paragraph in doc.paragraphs:
        pdf.multi_cell(0, 10, paragraph.text)

    pdf.output(output_path)

def pdf_to_excel(input_path, output_path):
    """Convert PDF to Excel."""
    reader = PdfReader(input_path)
    wb = Workbook()
    ws = wb.active

    row = 1
    for page in reader.pages:
        text = page.extract_text().split("\n")
        for line in text:
            ws.cell(row=row, column=1, value=line)
            row += 1

    wb.save(output_path)

def excel_to_pdf(input_path, output_path):
    """Convert Excel to PDF."""
    import pandas as pd

    wb = load_workbook(input_path)
    sheet = wb.active

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for row in sheet.iter_rows(values_only=True):
        row_text = " | ".join(str(cell) if cell else "" for cell in row)
        pdf.multi_cell(0, 10, row_text)

    pdf.output(output_path)

def pdf_to_ppt(input_path, output_path):
    """Convert PDF to PowerPoint."""
    reader = PdfReader(input_path)
    prs = Presentation()

    for page in reader.pages:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        text = page.extract_text()

        textbox = slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
        frame = textbox.text_frame
        frame.text = text

    prs.save(output_path)

def ppt_to_pdf(input_path, output_path):
    """Convert PowerPoint to PDF."""
    prs = Presentation(input_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)

    for slide in prs.slides:
        pdf.add_page()
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        pdf.multi_cell(0, 10, text)

    pdf.output(output_path)
